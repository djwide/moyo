"""Type-specific document loaders to extract raw text."""

from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

from .validators import ValidationConfig


@dataclass
class LoadedDoc:
    src_key: str
    mime: str
    bytes_sha256: str
    size_bytes: int
    meta: Dict[str, Any] = field(default_factory=dict)
    text: str = ""


def _sha256(data: bytes) -> str:
    import hashlib
    return hashlib.sha256(data).hexdigest()


def load_text_from_bytes(key: str, data: bytes, mime: str, cfg: ValidationConfig) -> LoadedDoc:
    size = len(data)
    bytes_hash = _sha256(data)
    
    # Dispatch by mime or extension
    if mime.startswith("text/") or key.lower().endswith((".txt", ".md")):
        text = _decode_text_bytes(data)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "text"}, text)
    
    if mime in {"text/html", "application/xhtml+xml"} or key.lower().endswith((".html", ".htm")):
        text, meta = _load_html(data)
        return LoadedDoc(key, mime, bytes_hash, size, {**meta, "loader": "html"}, text)
    
    if mime == "application/pdf" or key.lower().endswith(".pdf"):
        text, meta = _load_pdf(data, cfg)
        return LoadedDoc(key, mime, bytes_hash, size, {**meta, "loader": "pdfminer"}, text)
    
    if key.lower().endswith(".docx"):
        text = _load_docx(data)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "docx"}, text)
    
    if key.lower().endswith(".pptx"):
        text = _load_pptx(data)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "pptx"}, text)
    
    if key.lower().endswith(".xlsx"):
        text = _load_xlsx(data, cfg)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "xlsx"}, text)
    
    if key.lower().endswith((".csv",)):
        text = _load_csv(data)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "csv"}, text)
    
    if key.lower().endswith((".jsonl",)):
        text = _load_jsonl(data)
        return LoadedDoc(key, mime, bytes_hash, size, {"loader": "jsonl"}, text)
    
    # archives handled at pipeline level (iterate entries)
    
    # Fallback: treat as binary we cannot parse
    return LoadedDoc(key, mime, bytes_hash, size, {"loader": "unknown"}, "")


def _decode_text_bytes(data: bytes) -> str:
    try:
        from charset_normalizer import from_bytes
        res = from_bytes(data)
        best = res.best()
        if not best:
            return data.decode("utf-8", errors="strict")
        return str(best)
    except Exception:
        return data.decode("utf-8", errors="ignore")


def _load_html(data: bytes) -> tuple[str, Dict[str, Any]]:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(data, "lxml")
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        # Basic boilerplate removal: keep text in body
        body = soup.body.get_text("\n") if soup.body else soup.get_text("\n")
        text = f"{title}\n\n{body}".strip()
        return text, {"title": title}
    except Exception:
        return _decode_text_bytes(data), {}


def _load_pdf(data: bytes, cfg: ValidationConfig) -> tuple[str, Dict[str, Any]]:
    try:
        from pdfminer.high_level import extract_text_to_fp
        output = io.StringIO()
        extract_text_to_fp(io.BytesIO(data), output)
        return output.getvalue(), {}
    except Exception:
        # OCR hook placeholder (disabled by default)
        return "", {"warning": "pdf_parse_failed"}


def _load_docx(data: bytes) -> str:
    try:
        import docx  # python-docx
        f = io.BytesIO(data)
        d = docx.Document(f)
        parts = [p.text for p in d.paragraphs]
        return "\n".join([p for p in parts if p])
    except Exception:
        return ""


def _load_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = str(getattr(shape, "text"))
                    if t:
                        texts.append(t)
        return "\n".join(texts)
    except Exception:
        return ""


def _load_xlsx(data: bytes, cfg: ValidationConfig) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        texts: List[str] = []
        max_rows = cfg.max_rows_xlsx
        for ws in wb.worksheets:
            rows = 0
            for row in ws.iter_rows(values_only=True):
                rows += 1
                if rows > max_rows:
                    break
                str_cells = [str(c) for c in row if c is not None]
                if str_cells:
                    texts.append(",".join(str_cells))
        return "\n".join(texts)
    except Exception:
        return ""


def _load_csv(data: bytes) -> str:
    try:
        sio = io.StringIO(_decode_text_bytes(data))
        reader = csv.reader(sio)
        lines: List[str] = []
        for row in reader:
            lines.append(",".join(row))
        return "\n".join(lines)
    except Exception:
        return _decode_text_bytes(data)


def _load_jsonl(data: bytes) -> str:
    try:
        sio = io.StringIO(_decode_text_bytes(data))
        lines = []
        for line in sio:
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
                lines.append(json.dumps(obj, ensure_ascii=False))
            except Exception:
                lines.append(line)
        return "\n".join(lines)
    except Exception:
        return _decode_text_bytes(data)


